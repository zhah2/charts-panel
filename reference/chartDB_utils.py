# -*- coding: utf-8 -*-
"""
Created on Tue Jul 30 12:46:02 2024

@author: luoqin4
"""
import sqlalchemy as sa
import pandas as pd
import pyodbc
import re
from datetime import datetime

import requests
import os
import sys
from pathlib import Path
import warnings
import subprocess
from PIL import Image, ImageGrab
from io import BytesIO
import random, string
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import xlwings as xw
import win32com.client
from win32com.client import constants
import pyautogui
import win32gui
import win32con
import time
import shutil
import ctypes
import pythoncom
import gc
import logging




import json

# Global variable for the file path of opened_files_meta
OPENED_FILES_META_PATH = r"G:\Group\PORTFOLI\ASSETALLOC\Research\PythonProject\ResearchDB\ProductionCode\opened_file_meta.json"  # Update this path as needed






server_name = 'TDAMQ-SQLMD-PFO,3341'
db_name = 'ASSETALLOC'
driver = 'ODBC Driver 17 for SQL Server'
    


# Dictionary specifying the chartDB database's structure, i.e. column names for corresponding tables.
table_cols = {'chartDB_charts': ['ID', 'Name', 'Source', 'Location', 'Description', 'Status'],
              'chartDB_labels': ['Name', 'Type'],
              'chartDB_mapping': ['ID', 'Label']
              }



def get_engine(server_name, db_name, driver):
    '''
    Create and return the database engine object. Note 'AUTOCOMMIT' is turned on.
    '''
    engine = sa.create_engine(f'mssql+pyodbc://{server_name}/{db_name}?driver={driver}', isolation_level='AUTOCOMMIT')

    return engine
    
    
    

def get_table_as_df(conn, table_name):
    '''
    Get the specified table of a SQL Server database and return it as a dataframe.
    '''    
    query = f'SELECT * FROM {table_name}' # Create the query to pull out the table.
    
    df = pd.read_sql(query, con=conn) # Dump the table into a dataframe.
    
    return df
    


def get_curr_database(engine, results_folder, aggregate_labels=True, save_as_excel=True):
    
    '''
    Get a dump of the current database to specified folder.
    The result will be an Excel file with 2 tabs: one for charts and mapping, the other for existing labels.
    The optional argument labels_as_list controls whether labels will be aggregated into a lis
    
    '''     
    
    with engine.connect() as conn:
            
        df_mapping = get_table_as_df(conn, 'chartDB_mapping')
        df_charts = get_table_as_df(conn, 'chartDB_charts')
        df_labels = get_table_as_df(conn, 'chartDB_labels')
        
    if aggregate_labels == True:
        
        
        df_mapping_agg = df_mapping.groupby('ID')['Label'].agg(tuple) # Tuple is hashable so drop_duplicates() could be used later if needed.
        

        df_mapping_agg = df_mapping_agg.apply(lambda x: ', '.join(x)) # Convert tuple of labels to comma delimited string.
    
    df = df_charts.merge(df_mapping_agg, how='left', on=['ID']) # Add label info to charts dataframe to get a comprehensive df.
    
    # Insert the formatted label column to after 'Name' column, i.e., column number 2.
    label_col = df.pop('Label')
    df.insert(2, 'Label', label_col)
    
    if save_as_excel == True:
    
        # Generate file name with format: database_dump_YYYYMMDD_HHMM.xlsx
        ts = datetime.now().strftime('_%Y%m%d_%H%M')    
        file_name = 'database_dump' + ts + '.xlsx'
        file_path = results_folder + '\\' + file_name
       
        
        with pd.ExcelWriter(file_path, engine='xlsxwriter') as writer:
            
            df.to_excel(writer, sheet_name='charts', index=False)
            df_labels.to_excel(writer, sheet_name='labels', index=False)
            
            print(f'File {file_name} saved in {results_folder}.\n')

    return df, df_labels



def insert_chart(engine, input_dict=None, input_df=None):
    '''
    Takes either the standard user input dictionary for one single chart or a dataframe containing multiple charts and insert into database.
    
    The input_dict keys or the input_df column names should cover all the columns in table chartDB_charts and chartDB_mapping. 
    
    The insertion updates two tables: chartDB_charts and chartDB_mapping.
    
    If there're new labels detected from user input, the new labels will be inserted into chartDB_labels.
    '''
    
    if (input_dict is None) and (input_df is None):
        
        print('No chart data provided, please check.')
        
    elif input_dict is not None: # Get data from input dictionary. 
        
        # Get data for chartDB_charts as a dictionary.
        chart_dict = {k: input_dict[k] for k in table_cols['chartDB_charts']} 
        df_charts = pd.DataFrame(chart_dict, index = [0]) # Convert to dataframe, set ID to index.      
        
        # Get data for chartDB_mapping as a dictionary.
        if input_dict['Label'] is not None:
            mapping_dict = {k: input_dict[k] for k in table_cols['chartDB_mapping']} # Match keys with table columns.
            df_mapping = pd.DataFrame(mapping_dict).explode(column='Label') # Convert to dataframe.

    
    elif input_df is not None:  #Get input from provided dataframes   
       
        df_charts = input_df[table_cols['chartDB_charts']].drop_duplicates()
        df_mapping = input_df[table_cols['chartDB_mapping']].drop_duplicates()
        # If labels are in one line as comma separated strings, conver them to list. Otherwise this won't change anything.
        if df_mapping['Label'].notna().any(): # Do splitting only when there're valid strings, otherwise, the below will throw error because all labels are NA's or empty.
            
            df_mapping['Label'] = df_mapping['Label'].str.split(r',\s*') 
        
        df_mapping=df_mapping.explode(column='Label').drop_duplicates()

    # Check id format and keep only good ones.
    df_charts = df_charts[df_charts['ID'].apply(check_id_format)]
    df_mapping = df_mapping[df_mapping['ID'].isin(df_charts['ID'])].dropna(subset=['Label']) # Drop nan as well.
    
    # Check id existency. If already exist, notify user and drop the row.
    df_charts_dup, df_charts_new = check_duplicates(engine, df_charts, table_name='chartDB_charts', p_keys=['ID'], schema='dbo', notice=False)
    
    if not df_charts_dup.empty: # Notify user about already-exist charts that will not be inserted again.
        
        print(f'The following {len(df_charts_dup)} charts already exist in chartDB_charts. They will not be inserted.\n', df_charts_dup)         
               
    if not df_charts_new.empty: # Perform insertion for brand new charts in user input.
        
        # Drop duplicates and take only the new ID's.
        df_charts = df_charts_new
        
        if df_mapping is not None:
            
            df_mapping = df_mapping[df_mapping['ID'].isin(df_charts['ID'])]
        
            # Check if there're new labels and insert based on user spec.
            user_declined_list = check_labels_and_create(engine, df_mapping['Label'])
            
            if user_declined_list: # If there're new labels that user declined to insert, drop them from mapping.            
                
                df_mapping = df_mapping[~df_mapping['Label'].isin(user_declined_list)]   

        # Update database.
        insert_df_to_db(engine, df_charts, table_name='chartDB_charts', p_keys=['ID']) 
        
        if not df_mapping.empty:
            
            insert_df_to_db(engine, df_mapping, table_name='chartDB_mapping', p_keys=['ID', 'Label'])
        
        else:
            print('Please note that no mapping is inserted into chartDB_mapping table.')

    return df_charts, df_mapping

def check_labels_and_create(engine, labels):
    '''
    Check if all provided labels exist in chartDB_labels table.        
    If not exist, ask user whether to create the label, i.e. insert it into the label table.
    Return the list of labels that user declined to create.    
    
    Parameter:
        labels: an iterator (list or pandas series/dataframe etc. including a number of labels. )
        
    ''' 
    user_declined_list = []
    
    for label in labels: 
        
        if pd.notna(label):
       
            count_label = check_record(engine, table_name='chartDB_labels', col_name='Name', record_value=label) # check_reord() will return 0 if not found, otherwise 1.
            
            if count_label == 0:
                # Ask user whether to insert the new label if the label doesn't exist yet.
                user_cmd = input(f'New label detected: {label}. If you want to insert the new label into the database, please enter "y":')
    
                if user_cmd == 'y':
                    
                    type_input = input(f'Please enter the Type for label {label}:')
                    
                    insert_labels(engine, label_list=[{'Name': f'{label}', 'Type': f'{type_input}'}])            
                # print(f'The new label {label} has been inserted to chartDB_labels table. Type set to None for now.')
                
                elif user_cmd != 'y':
                    
                    print(f'New label {label} NOT inserted.\n')
                    
                    # Append the declined to create label to a list for return value.
                    user_declined_list.append(label) 
                
    return user_declined_list




def check_id_format(chart_id, warning=True):
    '''
    Check if an input chart_id is a 10-digit alphanumeric string in uppercase.
    Return True if yes, otherwise return False.
    User has the option to turn off the warning message by the argument 'warning'
    
    '''
    result = bool(re.fullmatch(r'[A-Z0-9]{10}', chart_id))
    
    if warning is True and result is False:  
            
        print(f'The provided chart ID {chart_id} does not follow the 10-digit uppercase alphanumeric format.\n')
    
    return result
        
    
    
def check_record(engine, table_name, col_name, record_value, condition_col=None, condition_value=None):
    '''
    Return 1 if a record already exists in specified table and column, otherwise 0.
    '''
    # Generate query for the count existing cases.
    
    if condition_col == None:
        
        query = sa.text(f"SELECT COUNT(1) FROM {table_name} WHERE {col_name} = '{record_value}'")
    
    else:
        
        query = sa.text(f"SELECT COUNT(1) FROM {table_name} WHERE {col_name} = '{record_value}' AND {condition_col} = '{condition_value}'")
        
        
        
    # Execute the check in the database.       
    with engine.connect() as conn:        
        
        try: 
            # result is a tuple with 1st element indicating the existence.
            result = conn.execute(query).first()          

            return result[0]
        
        except Exception as error:
            print(error)
        


def insert_labels(engine, label_list=None, label_df=None):
    '''
    Takes either a label_list or a dataframe with label name and type informationl, and then insert the labels into the chartDB_labels table of the database.
    
    The label_list should be a list of dictionary following the below format.
        
        new_labels = [{'Name': 'Labor', 'Type': 'Topics'},
                      {'Name': 'Inflation', 'Type': 'Topics'},
                      {'Name': 'US', 'Type': 'Country'}
                      ]   
        
    '''
    if label_list is None and label_df is None:
        
        print(f'No label data is provided, please check.')
        
    elif label_list is not None:
        
        df_labels = pd.DataFrame(label_list)    
    
    elif label_df is not None: # Get data from input file.
        
        # Input file should have the column names of table chartDB_lables, i.e. 'Name' and 'Type'.
        df_labels = label_df[table_cols['chartDB_labels']]
    
    # Check if there're duplicates in the input data.
    duplicates = df_labels[df_labels.duplicated('Name', keep=False)]
    
    if not duplicates.empty:
        
        print('The following duplicates found in the input, only the first occurance will be inserted.\n', duplicates, '\n')        
    
    df_labels.drop_duplicates(subset='Name', keep='first')

    insert_df_to_db(engine, df_labels, table_name='chartDB_labels', p_keys=['Name'])


    

def replace_chartID(engine, replace_id_list):
    '''
    Takes in a list of dictionary in below format, and rename the existing id to the new_value.
        
    replace_id_list = [{'ID': 'DQBORE80DE', 'New_Value': 'DQBORE80DX'},
                       {'ID': '123456789X', 'New_Value': '123456789Y'}
                       ]  
    
    The update operation is on the primary key of chartDB_charts table, the new id will automatically reflect in chartDB_mapping which has 'ID' as a freign key to the 'ID' in chartDB_charts.
    '''    

    for update in replace_id_list:
        
        # Check new id format.        
        if check_id_format(update['New_Value']) is True:            
        
            # Generate query for the update.
            update_record(engine, 
                          table_name='chartDB_charts', 
                          update_col='ID', 
                          new_value=update['New_Value'],
                          condition_col='ID',
                          condition_value=update['ID']                      
                          )  



def update_chart_info(engine, chart_id, update_info_list):
    '''
    Update the chart info based on provided update info list in following format:
        
        update_info_list = {'Name': 'NFIB Survey Raised Wage vs Planning to Raise,,, test',
                            'Source': 'Workspace,,,, test',
                            }     
    '''
    # Check if chart ID exists in charts table.        
    count_id = check_record(engine, table_name='chartDB_charts', col_name='ID', record_value=chart_id)

    if count_id == 0:
        
        print(f'Chart {chart_id} does not exist, please check.\n')
    
    else:
        for k, v in update_info_list.items():
            
            update_record(engine=engine, 
                          table_name='chartDB_charts', 
                          update_col=k, 
                          new_value=v, 
                          condition_col='ID', 
                          condition_value=chart_id
                          )
        

def remove_charts(engine, remove_chart_id_list):
    '''
    Takes a list of chart id and remove the corresponding charts and mapping.
    If a provided id does not exist, print a warning msg for user to check.
    
    '''
    
    df = pd.DataFrame(remove_chart_id_list, columns=['ID']) # Convert input to dataframe for duplicate check.
    
    df_duplicates, df_new = check_duplicates(engine, df, table_name='chartDB_charts', p_keys=['ID'], schema='dbo', notice=False)
    
    if not df_new.empty:
        
        print(f'The following {len(df_new)} ID provided do not exist in table chartDB_table, please check: \n', df_new, '\n')
    
    if not df_duplicates.empty:
        
        remove_records(engine, table_name='chartDB_charts', col_name='ID', record_list=df_duplicates['ID'].to_list())
    
    
    



def replace_labels(engine, replace_label_list):
    '''
    Takes in a label_list which is a list of dictionary in below format, and rename the existing lable name to the new_name.
        
    rename_label_list = [{'Name': 'Labor', 'New_Name': 'Labor_1'},
                         {'Name': 'Inflation', 'New_Name': 'Inflation_0'}
                         ]   
    
    The update operation is on the primary key of chartDB_labels table, the new name will automatically reflect in chartDB_mapping which has 'Label' as a freign key to the 'Name' in chartDB_labels.
    '''    

    for update in replace_label_list:
        
        # Generate query for the update.
        update_record(engine, 
                      table_name='chartDB_labels', 
                      update_col='Name', 
                      new_value=update['New_Name'],
                      condition_col='Name', 
                      condition_value=update['Name']                      
                      )   
               

def insert_mapping(engine, new_mapping_list):
    '''
    Takes a list of dictionary of new mapping information for existing charts and insert into the chartDB_mapping table. The input format is as follows:
        
        new_mapping = [{'ID': 'DQBORE80DE', 'Label': ['Labor', 'Inflation', 'US']},
                       {'ID': '123456789X', 'Label': ['Labor', 'Inflation', 'US']}               
                       ]
    
    If there're non-existing chart ID's, notify the user and drop the chart ID's.
    If there're non-existing labels, ask user whether to insert or not, and drop the non-existing labels that user declined to insert.
    If a mapping already exist, return error msg with detailed duplicate info.
    
    '''    
    df = pd.DataFrame(new_mapping_list).explode(column='Label')
    
    df_id = df['ID'].drop_duplicates().to_frame()
    
    df_id_duplicates, df_id_new = check_duplicates(engine, df_id, table_name='chartDB_charts', p_keys=['ID'], schema='dbo', notice=False)
    
    if not df_id_new.empty:
        
        print(f'The following {len(df_id_new)} ID do not exist in chartDB_Charts table. Please insert the charts first.\n', df_id_new, '\n')
    
    if not df_id_duplicates.empty:
        
        df = df[df['ID'].isin(df_id_duplicates['ID'])] # Update input df to only include existing id's.        
        
        # Check if there're new labels and insert based on user spec.
        user_declined_list = check_labels_and_create(engine, df['Label'])
        
        if user_declined_list: # If there're new labels that user declined to insert, drop them from mapping.            
            
            df = df[~df['Label'].isin(user_declined_list)] 

        if not df.empty:
            insert_df_to_db(engine, df, table_name='chartDB_mapping', p_keys=['ID', 'Label'])
        
        else:
            print('Please note that no mapping is inserted into chartDB_mapping table.')   
    



def check_duplicates(engine, df, table_name, p_keys=[], schema='dbo', notice=True):
    '''
    Check an input dataframe against the primary key columns of a table in the database.
    If there're duplicates, notify user.
    Return the duplicates and also the unique values.
    
    '''    
    df_duplicates = pd.DataFrame()
    df_new = pd.DataFrame()
    
    with engine.connect() as conn: # Get the primary-key columns of the table in database.
        
        df_table = get_table_as_df(conn, table_name)[p_keys] 
        
    # Before merge, get the order of the input dataframe.
    df['Order'] = range(len(df))

    # Get a merge of the df and the database table on primary key columns with indicator turned on.
    df = df.merge(df_table, how='left', on=p_keys, indicator=True)

    try:
        # Notify user if there're duplicate records already existing in database.
        df_duplicates = df[df['_merge'] == 'both'][p_keys + ['Order']] # Get duplicate records.
        df_duplicates = df_duplicates.sort_values('Order').drop(columns='Order')

        if not df_duplicates.empty:
            
            if notice == True: 
                
                print(f'The following {len(df_duplicates)} records already exist in column {p_keys} of table {table_name}.\n', df_duplicates, '\n')
    
    except:
        pass
        
    df_new = df[df['_merge'] == 'left_only'].sort_values('Order').drop(columns=['_merge', 'Order'])

    return df_duplicates, df_new



def insert_df_to_db(engine, df, table_name, p_keys=[], schema='dbo', success_notice=True):
    '''
    Check duplicates of a df against the specified p_keys columns of a table.
    If records already exist in primary key columns of table, they'll be dropped and user will be notified.
    Insert the unique values in p_key columns to the database.
    Both df and database table should contain columns named to the primary keys.
    
    '''
    
    # Check if any record in the dataframe already exists in the database. Notify user if there're.
    df_duplicates, df_new = check_duplicates(engine, df, table_name, p_keys=p_keys, schema='dbo')
    
    # Insert the not-already-exist records.
    if not df_new.empty:
        
        try:
            # Append the dataframe to database.
            df_new.to_sql(table_name, engine, schema=schema, if_exists='append', index=False)
            
            if success_notice == True:
                
                print(f'The following {len(df_new)} records have been successfully inserted in table {table_name}:\n', df_new[p_keys], '\n')
            
            return df_new[p_keys]            
        
        # If primary key constraint violated by duplicate chart ID, print the error with duplicate key information.
        except Exception as error: 
            
            print(f'Duplicate key detected while inserting into table {table_name} with following details: \n', error.args[0].split(table_name)[1], '\n')
        
            return None
    else:
        return None

    

def remove_labels(engine, remove_label_list=[], table_name='chartDB_labels', col_name='Name'):
    '''
    Remove labels in a given list. 
    Warn the user about the impact of removing labels: all relevant mapping will be gone.
    
    '''

    df = pd.DataFrame(remove_label_list, columns=['Name']) # Convert input to dataframe for duplicate check.
    
    df_duplicates, df_new = check_duplicates(engine, df, table_name=table_name, p_keys=['Name'], schema='dbo', notice=False)
    
    if not df_new.empty:        
            
        print(f'The following {len(df_new)} ID provided do not exist in table chartDB_table, please check: \n', df_new, '\n')
    
    if not df_duplicates.empty:    
    
        # Provide a warning msg when removing labels from chartDB_labels.
        print(f'Are you sure to remove all labels in list {df_duplicates[col_name].to_list()}? Please note all existing mapping related to these labels will be removed as well.')
    
        user_cmd = input('If continue, type "y":')
    
        if user_cmd == 'y':
        
            remove_records(engine, table_name, col_name, record_list=df_duplicates[col_name].to_list())
        
        else:
        
            print('Label removal operation cancelled.')
    



def remove_records(engine, table_name, col_name, record_list):
    '''
    Remove records in a given list from a table.
    
    The record_list should correspond to values of column 'col_name' in the specified table by 'table_name'.
    
    '''    
   
    # Generate query for the remove.
    query = sa.text(f'DELETE FROM {table_name} WHERE {col_name} IN ({str(record_list)[1:-1]})')
       
    # Execute the remove in the database.       
    with engine.connect() as conn:
        
        try: 
            conn.execute(query)   
            
            print(f'The following {len(record_list)} {col_name} records have been successfully removed from table {table_name}:\n', record_list)
        except Exception as error:
            print(error)
            


def remove_mapping(engine, remove_mapping_list, table_name='chartDB_mapping', p_keys=['ID', 'Label']):
    '''
    Remove list of mappings specified by user in the following format:
        
        remove_mapping_list = [{'ID': 'DQBORE80DE', 'Label': ['Labor3', 'Inflation', 'US']},
                               {'ID': '123456789X', 'Label': ['Labor', 'Inflation', 'US']}  
                               ]
        
    The values in p_keys list should correspond to the keys in the removing_list as well as the primary keys in the database mapping table.
    
    '''           
    
    r_list = deepcopy(remove_mapping_list) # Get a deep copy of the remove list for the for loop.
    
    for mapping in r_list:
        
        # Get a copy of the original mapping for later removal operation.
        for original_mapping in remove_mapping_list:
            
            if original_mapping == mapping:
                
                org_mapping = original_mapping
        
        # Check if chart ID exists in charts table.        
        count_id = check_record(engine, table_name='chartDB_charts', col_name='ID', record_value=mapping['ID'])

        if count_id == 0:
            
            print(f'Chart {mapping["ID"]} does not exist, please check the ID in the remove_mapping_list.\n')
            
            remove_mapping_list.remove(mapping) # Drop the non-existing chart id from original remove list.
            
            continue
        
        # Check if any wrong label name specified by user.
        for label in mapping['Label']:

            count_label = check_record(engine, table_name='chartDB_labels', col_name='Name', record_value=label) # check_reord() will return 0 if not found, otherwise 1.

        
            if count_label == 0:
            
                print(f'Label "{label}" does not exist, please check the label name for chart {mapping["ID"]} in the remove_mapping_list.\n') 
                
                # Drop the non-existing label from the original remove list.
                org_mapping['Label'].remove(label) 
                
            else: # Further check if the mapping exists.
                            
                count_mapping = check_record(engine, table_name='chartDB_mapping', col_name='ID', record_value=mapping['ID'], condition_col='Label', condition_value=label) 
                
                
                if count_mapping == 0: # If mapping not exist, there's nothing to remove.
                
                    print(f'Mapping ({mapping["ID"]}, {label}) does not exist, please check.\n') 
                    
                    # Drop the non-existing label from the original remove list.
                    org_mapping['Label'].remove(label) 

        
        # If not all labels are dropped due to not existing, continue to remove the mapping.
        if org_mapping['Label']:             

            # Generate the query for removing the mapping.
            condition = ''
            
            for key in p_keys:
                
                new_condition = key + ' IN (' + str_if_list(mapping[key]) + ')'
                
                if condition == '':
                    
                    condition = new_condition
                    
                else:                    
                    condition = condition + ' AND ' + new_condition
                
            # Generate query for the remove.
            query = sa.text(f'DELETE FROM {table_name} WHERE {condition}') 
                
            # Execute the remove in the database.       
            with engine.connect() as conn:
                
                try: 
                    conn.execute(query)   
                    
                    print(f'Mapping ({mapping["ID"]}, {mapping["Label"]}) successfully removed.')
                    
                except Exception as error:
                    print(error)
               


def str_if_list(input_data):
    '''
    If input is a list, return the list content as strings delimited by comma.
    If input is string, return the string wrapped by ''
    Return error msg otherwise.
    '''
    if isinstance(input_data, list):
        
        return str(input_data)[1:-1]
    
    elif isinstance(input_data, str):
        
        return "'" + input_data + "'"
    
    else:
        print(f'Unexpected type of input for str_if_list: {type(input_data)}.\n')


            
def update_record(engine, table_name, update_col, new_value, condition_col, condition_value):
    '''
    Update a record in a table with given condition and new value.
    The column and value inputs correspond to the following query:
        
        UPDATE {table_name} SET {update_col} = '{new_value}' WHERE {condition_col} = '{condition_value}
        
    '''
    # Check if condition_value exists in specified condition column.        
    count = check_record(engine, table_name=table_name, col_name=condition_col, record_value=condition_value)

    if count == 0:
        
        print(f'{condition_col} {condition_value} does not exist in table {table_name}, please check.\n')
        
        return  
    
    else:
    
        # Generate query for the update.
        query = sa.text(f"UPDATE {table_name} SET {update_col} = '{new_value}' WHERE {condition_col} = '{condition_value}'")
        
        
        # Execute the remove in the database.
           
        with engine.connect() as conn:
            
            try: 
                conn.execute(query)   
                
                print(f'{update_col} has been successfully updated to {new_value} where {condition_col} = {condition_value} in table {table_name}.\n')
                
            except Exception as error: 
                
                print(f'Duplicate primary key detected while update {update_col} to {new_value} in table {table_name} with following details: \n', error.args[0].split(table_name)[1], '\n')            
            
            


def get_chart_df_by_labels(engine, label_set_list):
    '''
    Get a dataframe with chart name location info based on a user provided label list.
    
    Each list item is a set of labels which means user would like to extract charts whose labels contain this set, i.e. logical 'AND' within the set. E.g., {'Inflation', 'US'} requires all charts whose labels contain 'Inflation' AND 'US'.
    
    The list items have a logical 'OR' relationship among them, e.g., [{'Inflation', 'US'}, {'Labor'}] should extract all charts either with labels 'Inflation' and 'US' or with label 'Labor'.
        
    '''
    
    with engine.connect() as conn:
        
        df_charts = get_table_as_df(conn, 'chartDB_charts').set_index('ID')
        df_labels = get_table_as_df(conn, 'chartDB_labels')
        df_mapping = get_table_as_df(conn, 'chartDB_mapping').set_index('ID')
    
    
    # Check if all labels in label_set_list are valid labels and drop the invalid ones.
    db_labels = set(df_labels['Name']) # Get the set of labels in database.
    
    for s in list(label_set_list): # Iterate over a copy of the list to avoid modifying it during iteration. Note s is reference to original list items.
        
        for label in list(s): # Iterate over a copy of the set to avoid modifying it during iteration.
        
            if label not in db_labels:
                
                print(f'Label "{label}" does not exist in table chartDB_labels, please check. This label will NOT be used for getting charts.\n')
                s.remove(label) # Drop the non-existing label from the original set in the list.                
                
        # If original set has been reduced to an empty set, drop it from original set list. 
        # Including empty set will lead to all charts been grabbed since empty set is a subset of any set.
        if not s: 
            
            label_set_list.remove(s)                    
        
    # Get a dataframe with 'Label' column flattened to tuple for each chart.
    df_flat = df_mapping.groupby('ID')['Label'].agg(tuple) # Tuple is hashable so drop_duplicates() could be used later.
    df_charts = pd.concat([df_charts, df_flat], axis=1) # Add label info to charts dataframe.
    
    df = pd.DataFrame() # Initialize the output dataframe.
    
    for label_set in label_set_list:
                
        # Select charts contains the current set of labels that is being looped on.
        t = df_flat.map(lambda x: label_set.issubset(x))
        
        df = pd.concat([df, df_charts[t]])
    
    df = df.reset_index().drop_duplicates()
            
    return df




def get_chart_df_by_ids(engine, id_list):
    '''
    Get a dataframe with chart name, label and location info based on a user provided chart id list.
    The returned df should have the following columns: 
    '''
    
    # check id existency.
    df_input = pd.DataFrame(id_list, columns=['ID'])
    df_exist, df_new = check_duplicates(engine, df_input, table_name='chartDB_charts', p_keys=['ID'], schema='dbo', notice=False)

    
    df = pd.DataFrame() # Initialize the output dataframe.
    
    if not df_new.empty:
        
        print(f"The folloiwng ID's don't exist in database, please check.", df_new, '\n')
    
    if not df_exist.empty: 
        
        # Get mapping and chart info for existing chart ids.
        with engine.connect() as conn:
            
            df_mapping = get_table_as_df(conn, 'chartDB_mapping')
            df_charts = get_table_as_df(conn, 'chartDB_charts')
            
        
        df_flat = df_mapping.groupby('ID')['Label'].agg(tuple) # Tuple is hashable so drop_duplicates() could be used later.
        
        # Get a merge of the df and the database table on primary key columns with indicator turned on.
        df_exist['Order'] = range(len(df_exist))
        df = df_exist.merge(df_charts, how='left', on=['ID'])

        df = df.merge(df_flat, how='left', on=['ID']).sort_values('Order').drop(columns='Order')
        
   
        
     
    else:
        print(f'No valid charts extracted from database, please check.')        
        
    return df
    

def get_charts_and_generate_ppt(engine, results_folder, ppt_file_name='Research Package', label_set_list=None, chart_id_list=None, chart_df=None, rerun_user_chart=False, user_name=True, time_stamp=True, cover=None, cover_page_title=None, add_chart_meta=False, add_page_num=False, save_pdf=False):
    '''
    Generate powerpoint file based on either label or id inputs. 
    The file name is default to 'charts' if not specified.
    Timestamp will be added to file name if time_stamp is True.
    If cover slide is provide, the first slide will be the given cover slide.
    If rerun_user_chart is True, the TDAM Model charts will be refreshed before inserting into the pptx file.
    If save_pdf is True, a PDF copy of the generated PowerPoint file will be saved.
    '''
    if label_set_list is None and chart_id_list is None and chart_df is None:
        
        print('No input provided, please specify either a list of label sets or a list of chart ids or a dataframe with chart ids and titles.')
        
    elif label_set_list is not None:
        
        df = get_chart_df_by_labels(engine, label_set_list)
    
    elif chart_id_list is not None:
        
        df = get_chart_df_by_ids(engine, chart_id_list)
    
    elif chart_df is not None:
        
        df = get_chart_df_by_ids(engine, chart_df['ID'].tolist())
        
        df = df.drop_duplicates()
        

        chart_df = chart_df[chart_df['ID'].isin(df['ID'])]
        
        chart_df['Order'] = range(len(chart_df))

        
        if 'Slide Title' in chart_df.columns:
            

            df = pd.merge(chart_df[['ID', 'Slide Title', 'Slide Number', 'Position', 'Position_Adj','Strech Factor', 'Order']], df, on='ID', how='left').sort_values('Order').drop(columns='Order')

        

        
    if not df.empty:
        
        generate_ppt_file(df, results_folder, 
                          ppt_file_name=ppt_file_name, 
                          rerun_user_chart=rerun_user_chart, 
                          user_name=user_name, 
                          time_stamp=time_stamp, 
                          cover=cover, 
                          cover_page_title=cover_page_title,
                          add_chart_meta=add_chart_meta,
                          add_page_num=add_page_num,
                          save_pdf=save_pdf
                          )
        
    
        return df
    
    else:
                
        print('Slides are not generated.')
    
    
    
    
    

def generate_ppt_file(df, results_folder, archive_chart_imgs=False, ppt_file_name='charts', rerun_user_chart=True, user_name=True, time_stamp=True, cover=None, cover_page_title=None, add_chart_meta=False, add_page_num=False, save_pdf=False):
    '''
    Takes an input df with the following chart info: id, source, location.
    Generates a pptx file containing the charts in the input dataframe.
    Note different source would need different extraction methods:
        
        'Workspace' usually has an URL.
        'TDAM Model' would have the chart png file path.
        etc.
        
    '''       
    # Get layout numbers from cover page's slide master.
    blank_layout = 1 
    regular_layout = 2
    part_header_layout = 3
    left_right_layout = 4 
    top_btm_layout = 5
    quadrant_layout = 6
    

    if user_name:
        
        file_name = ppt_file_name + '_' + os.environ.get('USERNAME') 
        
    else:
        file_name = ppt_file_name
    
    if time_stamp:
        
        ts = datetime.now().strftime('_%Y%m%d_%H%M')
        
        img_folder_path = results_folder + '\\' + file_name + ts + '_images' # Generate the img archive folder name.
        
        file_name = file_name + ts + '.pptx' # Generate the file name        
        
    else:
        
        img_folder_path = results_folder + '\\' + file_name + '_images'
        
        file_name = ppt_file_name + '.pptx'
    
    os.makedirs(img_folder_path, exist_ok=True) # Create folder to archive chart images in results folder.
    

    
    if cover: # Use the given cover slide if provided.
                
        prs = Presentation(cover)
        
        cover = prs.slides[0] # Provide cover template should have only one slide which is the cover page.
        
        # Set file name to cover title. Requires cover title to be the 1st shape in the cover template.
        cover.shapes[0].text_frame.text = cover_page_title 
            
        date_str = datetime.now().strftime('%B %d, %Y') 
        
        cover.shapes[1].text_frame.text = date_str # Require date frame to be the 2nd shape in the cover template.   
        
        page_num = 2 # Initialize starting page number.
        
        
    else: # If cover is not specified, create a blank presentation.
        prs = Presentation() 
        
        page_num = 1 # Initialize starting page number.
                    
     
    
    # If the chart df doesn't have 'Slide Title' Column, create one and default the titles to chart names.
    if 'Slide Title' not in df.columns:
        
        df['Slide Title'] = df['Name']
        
        
    # Initialize current_slide and page_num.    
    current_slide = None
    new_page_flag = False
    
    
    current_slide_number = None
    chart_meta_textbox = None
    accumulated_meta_info = ""  # Variable to accumulate meta info for each slide
    
    # Default left, top, width and height numbers
    L = 0.28
    T = 0.7 
    W = 9.4 
    H = 6.5
    Gap = 0.05
    default_strech_factor = 0.4

    for index, row in df.iterrows():
        
        # Get image strech factor.
        if pd.isna(row['Strech Factor']):
            
            strech_factor = default_strech_factor
        
        else:
            
            strech_factor = row['Strech Factor']               
        

        # Get position adjustments.
        if pd.isna(row['Position_Adj']):
            
            l_adj, t_adj, w_adj, h_adj = [0, 0, 0, 0]
        
        else:
            
            l_adj, t_adj, w_adj, h_adj = [Inches(float(i)) for i in row['Position_Adj'].split(',')]
        
        
        # Process part header slide first.
        if row['ID'] == 'PARTHEADER':           
            
            
            new_page_flag = True # Turn on new_page_flag.
            accumulated_meta_info = ""
            
            if cover:
                
                slide = prs.slides.add_slide(prs.slide_layouts[part_header_layout])
                
                slide.shapes[0].text_frame.text = str(row['Slide Title']) # Set file name to cover title.
                    
            
            else:
                # Add a slide with blank layout.
                slide = prs.slides.add_slide(prs.slide_layouts[blank_layout])
            
                # Add a text box in the middle of the slide.
                part_header_box = slide.shapes.add_textbox(left=Inches(1) + l_adj, top=Inches(4) + t_adj, width=Inches(8) + w_adj, height=Inches(1) + h_adj)
                
                part_header_box.text = str(row['Slide Title'])                  
                
                title_box.text_frame.paragraphs[0].font.size = Pt(24)
                title_box.text_frame.paragraphs[0].font.bold = True
                title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        
        
        
        
        # Add Alternatives Slides.
        elif row['ID'] == 'ALTSSLIDES':
            
            # If need to take new screenshots, run below, otherwise, take existing ones in folder.
            if rerun_user_chart == True:
            
                # Load the .pptx file from the location folder (assuming only one exists)
                pptx_file_name = next((f for f in os.listdir(row['Location']) if f.endswith('.pptx')), None)
                
                if pptx_file_name:
                    
                    pptx_file_path = os.path.join(row['Location'], pptx_file_name)              
                    
                    
                    
                    slide_screenshots_dict = get_slides_screenshots(pptx_file_path, crop=False, screenshot_area=None, hide_page_numbers=True, wait_time=5, desired_dpi=120)                
                    
                    ext_prs = Presentation(pptx_file_path)  # Load the .pptx file found
                    
                    prs, page_num = add_external_slide_screenshots(curr_prs=prs, ext_prs=ext_prs, current_page_number=page_num, slide_screenshots_dict=slide_screenshots_dict, position=(0,0, 10, 7.5), add_page_num=add_page_num)      #position=(Inches(0.28), Inches(0.7), Inches(9.4), Inches(6.5))
                    
                                    
                else:
                    print(f"No .pptx file found in {row['Location']}.")
                    #raise FileNotFoundError("No .pptx file found in the specified folder.")  
                
                
        
        # Add TDAM Equity Team Earnings Revision Slides.
        elif row['ID'] == 'TDEQSLIDES':
            
            # If need to take new screenshots, run below, otherwise, take existing ones in folder.
            if rerun_user_chart == True:
            
                # Load the .pptx file from the location folder (assuming only one exists)
                pptx_file_name = next((f for f in os.listdir(row['Location']) if f.endswith('.pptx')), None)
                
                if pptx_file_name:
                    
                    pptx_file_path = os.path.join(row['Location'], pptx_file_name)   
                    
                    
                    slide_screenshots_dict = get_slides_screenshots(pptx_file_path, screenshot_area={'Default': (0.4, 1.2, 9.6, 5.9)}, hide_page_numbers=True, wait_time=5, desired_dpi=120)          
                                   
    
                                    
                else:
                    print(f"No .pptx file found in {row['Location']}.")
                    #raise FileNotFoundError("No .pptx file found in the specified folder.") 
        
          
                            
        else:
            if current_slide is None or row['Slide Number'] != current_slide_number:
                
                new_page_flag = True # Turn on new_page_flag.
                accumulated_meta_info = ""

                # Add slide title if not PARTHEADER
                if cover:
                    # Use the regular_layout of the cover
                    slide = prs.slides.add_slide(prs.slide_layouts[regular_layout])                    
                    title_shape = slide.shapes.title
                    if title_shape:
                        title_shape.text = row['Slide Title']
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[blank_layout])   # Assuming the first layout is the regular layout
                    # Create a title box at the top
                    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                    text_frame = title_shape.text_frame
                    text_frame.text = row['Slide Title']
                    
                current_slide = slide
                current_slide_number = row['Slide Number']

                # Add chart meta info
                if add_chart_meta:
                    
                    # Create a text box at the bottom of the slide
                    chart_meta_textbox = slide.shapes.add_textbox(Inches(L), Inches(T + H), Inches(W), Inches(0.8))
                    chart_meta_textbox.text_frame.clear()
                    
                    

            img_path = get_chart_img(row['Source'], row['Location'], row['ID'], img_archive_folder=img_folder_path, rerun_user_chart=rerun_user_chart)
            
            if img_path:
                
                
                left, top, width, height = Inches(L), Inches(T), Inches(W), Inches(H)
                
                if row['Position'] == 'Center':
                    pass
                
                elif row['Position'] == 'Left':
                    width, height = Inches(W/2), Inches(H)
                    
                elif row['Position'] == 'Right':
                    left, width, height = Inches(L + W/2 + Gap), Inches(W/2), Inches(H)
                    
                elif row['Position'] == 'Top':
                    height = Inches(H/2)
                    
                elif row['Position'] == 'Bottom':
                    top, height = Inches(T + H/2), Inches(H/2)
                    
                elif row['Position'] == 'Top Left':
                    width, height = Inches(W/2), Inches(H/2)
                    
                elif row['Position'] == 'Top Right':
                    left, width, height = Inches(L + W/2 + Gap), Inches(W/2), Inches(H/2)
                    
                elif row['Position'] == 'Bottom Left':
                    top, width, height = Inches(T + H/2), Inches(W/2), Inches(H/2)
                    
                elif row['Position'] == 'Bottom Right':
                    left, top, width, height = Inches(L + W/2 + Gap), Inches(T + H/2), Inches(W/2), Inches(H/2)
                    
                left, top, width, height = left + l_adj, top + t_adj, width + w_adj, height + h_adj

                # Add the picture without specifying dimensions
                picture = slide.shapes.add_picture(img_path, left, top)
                

                # Calculate scaling factor to fit within the specified area while maintaining aspect ratio
                scale_width = width / picture.width
                scale_height = height / picture.height
                scale = min(scale_width, scale_height)
                strech_scale = scale + strech_factor * abs(scale_width - scale_height)
                
                # Resize the picture while maintaining aspect ratio
                if scale_width <= scale_height:
                    picture.width = int(picture.width * scale)
                    picture.height = int(picture.height * strech_scale)
                
                else:
                    picture.width = int(picture.width * strech_scale)
                    picture.height = int(picture.height * scale)
                
                # Center the picture within the specified area
                picture.left = int(left + (width - picture.width) / 2)
                picture.top = int(top + (height - picture.height) / 2)
                
                #print(row['ID'], picture.left, picture.top)
                
                # Accumulate meta info if add_chart_meta is True
                if add_chart_meta and chart_meta_textbox:
                    
                    
                    if pd.isna(row['Label']):
                        new_meta_info = f"ID: {row['ID']};  Source: {row['Source']}"
                    else:
                        new_meta_info = f"ID: {row['ID']};  Source: {row['Source']};  Label: {', '.join(row['Label'])}"
                        
                    if accumulated_meta_info:
                        accumulated_meta_info += "\n" + new_meta_info
                    else:
                        accumulated_meta_info = new_meta_info
                        
                    # Update the chart meta text box
                    text_frame = chart_meta_textbox.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = accumulated_meta_info.strip()  # Remove leading/trailing whitespace
                    p.font.size = Pt(8)
                

        if new_page_flag:
            
            if add_page_num:
            
                # Add page number at bottom right corner.
                page_num_text_box = slide.shapes.add_textbox(left=Inches(9.6), top=Inches(6.8), width=Inches(0.5), height=Inches(0.5))
                page_num_text_frame=page_num_text_box.text_frame
                # Add chart ID at the bottom
                page_num_para = page_num_text_frame.add_paragraph()
                page_num_para.text = str(page_num)
                page_num_para.font.size = Pt(10)
                page_num_para.font.color.rgb = RGBColor(0, 0, 0)  # Black color 
                               
            print(f'Slide {page_num} created.\n')
            logging.info(f'Slide {page_num} created.')
            page_num +=1
        
        
        new_page_flag = False # Reset new_page_flag false by the end of each loop.
        
        
    ppt_file_path = results_folder + '\\' + file_name
    
    prs.save(ppt_file_path) 
    print(f'{file_name} saved in {results_folder}.')  
    logging.info(f'{file_name} saved in {results_folder}.')
    
    
    
    #### Turn page number on if it's not. And save as pdf if needed.
    
    # Wait 3 seconds before opening the generated file to turn on page number.
    time.sleep(3)
        
    powerpoint = win32com.client.Dispatch('PowerPoint.Application')
    #powerpoint = win32com.client.gencache.EnsureDispatch('PowerPoint.Application')
    presentation = powerpoint.Presentations.Open(ppt_file_path)
    
    # Wait for the file to fully open
    time.sleep(3)

    for slide_index, slide in enumerate(presentation.Slides, start=0):
        
        if slide_index == 0:
            
            slide.HeadersFooters.SlideNumber.Visible = False
        
        else: 
            
            slide.HeadersFooters.SlideNumber.Visible = True            
    
    # Save PDF copy if requested
    if save_pdf:
        try:
            
            # Generate PDF file path
            pdf_file_path = ppt_file_path.replace('.pptx', '.pdf')
            
            # Save as PDF
            presentation.SaveAs(pdf_file_path, 32)  # 32 is the PDF format code
            
            print(f'PDF copy saved as: {pdf_file_path}')
            logging.info(f'PDF copy saved as: {pdf_file_path}')
            
        except Exception as e:
            print(f'Error saving PDF copy: {str(e)}')
            logging.error(f'Error saving PDF copy: {str(e)}')
    
    
    # Save and close.
    presentation.Save()
    
    time.sleep(3)
    try:
        presentation.Close()

    
    except Exception as e:
        print("An error has occurred while closing the generated ppt file:", e)
    
    
    
    # Clean up temporary folder
    if not archive_chart_imgs:
        
        try:
            
            shutil.rmtree(img_folder_path)
            print('Achived chart images removed from results folder.')
        
        except:
            
            print('Error occured when removing temporary image archive folder.')

    


def get_slides_screenshots(file_path, crop=True, screenshot_area=None, hide_page_numbers=False, wait_time=5, desired_dpi=300):
    '''
    screenshot_area is a dictionary that specifies the screenshot area for each slide.
    The keys are the slide numbers, and the values are tuples that specify the screenshot area in inches.
    example:
    screenshot_area = {
        1: (0, 0, 10, 7.5),  # Tuple for slide 1
        2: (1, 1, 8, 6),     # Tuple for slide 2
        'Default': (0, 0, 8, 6)  # Default tuple for other slides
    }
    '''

    # Step 1: Make a copy of the original file
    copy_path = file_path.replace('.pptx', '_copy.pptx')
    shutil.copy(file_path, copy_path)  # Copying the file instead of renaming

    slide_screenshots_dict = {}
    try:
        # Step 2: Open the copy
        powerpoint = win32com.client.Dispatch('PowerPoint.Application')
        #powerpoint = win32com.client.gencache.EnsureDispatch('PowerPoint.Application')
        presentation = powerpoint.Presentations.Open(copy_path)
        
        # Wait for the file to fully open
        time.sleep(wait_time)
        
        
        # Get the full size of the slide
        prs = Presentation(copy_path) # Load the presentation to prs object
        emu_per_inch = 914400
        full_width = prs.slide_width/emu_per_inch
        full_height = prs.slide_height/emu_per_inch
        

        # Step 3: Take screenshots of each slide
        for slide_index, slide in enumerate(presentation.Slides, start=0):
            if hide_page_numbers:
                slide.HeadersFooters.SlideNumber.Visible = False
                
                
            # Export the full slide as a PNG with desired DPI.
            export_width = int(full_width * desired_dpi)  # Convert width from inches to pixels
            export_height = int(full_height * desired_dpi)  # Convert height from inches to pixels
            
            full_slide_path = os.path.join(os.path.dirname(file_path), f"{slide_index}_full.png")
            slide.Export(full_slide_path, "PNG", export_width, export_height)
            
            
           

            # Crop the image based on the specified area
            with Image.open(full_slide_path) as img:
                
                if crop == False: # If no need to crop, save the full screenshot
                    cropped_img = img
                
                else: # If need to crop, determine the screenshot area and crop the image accordingly.
                
                    # Determine the screenshot area
                    if screenshot_area is None:
                        # Default to full slide size
                        left, top, width, height = 0, 0, full_width, full_height
                        #left, top, width, height = 0, 0, slide.Master.Width, slide.Master.Height
                    elif slide_index in screenshot_area:
                        # Use the user-specified screenshot area directly in inches
                        left, top, width, height = screenshot_area[slide_index]
                    else:
                        default_area = screenshot_area['Default'] if screenshot_area and 'Default' in screenshot_area else (0, 0, full_width, full_height)  # Default values in inches
                        # default_area = screenshot_area['Default'] if screenshot_area and 'Default' in screenshot_area else (0, 0, slide.Master.Width, slide.Master.Height)  # Default values in inches
                        left, top, width, height = default_area
                    
                    
                    
                    # Get the image dimensions in pixels.
                    img_width, img_height = img.size
    
                    # Get the DPI from the image
                    dpi = img.info.get('dpi', (desired_dpi, desired_dpi))  # Default to desired DPI if not found
                    dpi_x, dpi_y = dpi  # Unpack the DPI values
    
                    # Convert inches to pixels using the desired DPI
                    left_px = int(left * desired_dpi)  # Convert left from inches to pixels
                    top_px = int(top * desired_dpi)    # Convert top from inches to pixels
                    width_px = int(width * desired_dpi)  # Convert width from inches to pixels
                    height_px = int(height * desired_dpi)  # Convert height from inches to pixels
    
                    # Ensure cropping dimensions are within image bounds
                    left_px = max(0, left_px)
                    top_px = max(0, top_px)
                    width_px = min(width_px, img_width - left_px)
                    height_px = min(height_px, img_height - top_px)
    
                    # Crop the image
                    cropped_img = img.crop((left_px, top_px, left_px + width_px, top_px + height_px))
                    
                    
                cropped_img.save(os.path.join(os.path.dirname(file_path), f"{slide_index}.png"))  # Save the cropped image

            slide_screenshots_dict[slide_index] = os.path.join(os.path.dirname(file_path), f"{slide_index}.png")
        
        
        # Step 4: Close the file without saving changes
        time.sleep(3)
        try:
            presentation.Close()  # Use SaveChanges=False to not save changes
        except Exception as e:
            print("An error has occurred while closing the presentation:", e)

    except Exception as e:
        print("An error has occurred while taking screenshots:", e)

    # Step 5: Wait for 3 seconds before removing the copy
    time.sleep(3)
    try:
        os.remove(copy_path)  # Remove the copied file
    except Exception as e:
        print("An error has occurred while removing the copied file:", e)
        
    # Return the slide number and corresponding screenshot file paths
    return slide_screenshots_dict
    #return pd.DataFrame(list(slide_screenshots_dict.items()), columns=['Slide Number', 'Screenshot Path'])


        
                
def add_external_slide_screenshots(curr_prs, ext_prs, current_page_number, slide_screenshots_dict, position=(0, 0, 10, 7.5), add_page_num=False):
    
   
    # Loop through each slide in the external presentation
    for ext_slide_index, ext_slide in enumerate(ext_prs.slides):
        shapes = ext_slide.shapes

        if len(shapes) == 1:
            # If there's only one shape, create a new slide with 'part_header_layout'
            part_header_layout = curr_prs.slide_layouts[3]  # Assuming layout index 1 is 'part_header_layout'
            new_slide = curr_prs.slides.add_slide(part_header_layout)
            title_shape = new_slide.shapes.title
            title_shape.text = shapes[0].text  # Set the text content
            

        
        elif len(shapes) > 1:

            # If there are multiple shapes, create a new slide and add the screenshot
            new_slide = curr_prs.slides.add_slide(curr_prs.slide_layouts[1])  # Use a blank layout
            
            # Get the screenshot for the current slide index
            screenshot_path = slide_screenshots_dict.get(ext_slide_index)

            if screenshot_path:
                # Add the screenshot to the specified position with default dimensions
                left = Inches(position[0])
                top = Inches(position[1])
                width = Inches(position[2])
                height = Inches(position[3])
                
                picture = new_slide.shapes.add_picture(screenshot_path, left=left, top=top)            
           
                # Calculate scaling factor to fit within the specified area while maintaining aspect ratio
    
                scale_width = width / picture.width
                scale_height = height / picture.height
                scale = min(scale_width, scale_height)

    
                picture.width = int(picture.width * scale)
                picture.height = int(picture.height * scale)            

                
                # Center the picture within the specified area
                picture.left = int(left + (width - picture.width) / 2)
                picture.top = int(top + (height - picture.height) / 2)

            
        # Add page number at bottom right corner.
        if add_page_num:
            page_num_text_box = new_slide.shapes.add_textbox(left=Inches(9.6), top=Inches(6.8), width=Inches(0.5), height=Inches(0.5))
            page_num_text_frame=page_num_text_box.text_frame
            # Add chart ID at the bottom
            page_num_para = page_num_text_frame.add_paragraph()
            page_num_para.text = str(current_page_number)
            page_num_para.font.size = Pt(10)
            page_num_para.font.color.rgb = RGBColor(0, 0, 0)  # Black color 
                           
        print(f'Slide {current_page_number} created.')
        logging.info(f'Slide {current_page_number} created.')
        current_page_number +=1            


    # Return the updated current presentation and the last page number
    return curr_prs, current_page_number
                


    
      
    
    
def get_chart_img(source, location, chart_id, img_archive_folder, rerun_user_chart=True):
    '''
    Get the image of a chart based on source and location information.
    For source='TDAM Model', if 'refreash' flag is True, 
    '''
    
    if source == 'Workspace': # Location for 'Workspace' should be an URL.
               
        # Turn off InsecureRequestWarning for the usage of request package.
        warnings.simplefilter('ignore', requests.packages.urllib3.exceptions.InsecureRequestWarning)
        
        response = requests.get(location, stream=True, verify=False)
        
        img = Image.open(response.raw)
        img_archive_path = img_archive_folder + f'\\{chart_id}.png'
        img.save(img_archive_path)       
        
        return img_archive_path
    
    elif source == 'TDAM Model':
        
        
        if rerun_user_chart == True:
            
            user_folder = os.path.dirname(location) # Get user code folder path.
            
            for file_name in os.listdir(user_folder):
                
                file_path = os.path.join(user_folder, file_name)
                
                if os.path.isfile(file_path) and file_name.endswith('.py'):
                    
                    print(f'Start running {file_name} in folder [{user_folder}].\n')
                    logging.info(f'Start running {file_name} in folder [{user_folder}].\n')
                    
                    try:                                                
                        
                        result = subprocess.run(['python', file_path], capture_output=True, check=True,  text=True) 
                        
                        #if result.returncode == 0:
                             
                        print(f'{file_name} run successfully.\n')
                        logging.info(f'{file_name} run successfully.\n')

                        
                        
                    except subprocess.CalledProcessError as call_error:
                        
                        #print(f'Error runing {file_name}: ', call_error.stderr, '\n')
                        e_lines = call_error.stderr.strip().splitlines()
                        
                        if e_lines:
                            e_msg = call_error.stderr.strip().splitlines()[-1]                                
                            
                            print(f'{file_name} run failed with error: ', e_msg, '\n')
                            
                            logging.error(f'{file_name} run failed with error: {e_msg}')
                            
                            
                        
                        #sys.exit('The above error is critical, run stopped 1.')                        
                        continue
                        
                    except Exception as err:
                        
                        print(f'Unexpected error happend when calling {file_name}: ', err, '\n')
                        
                        logging.error(f'Unexpected error happend when calling {file_name}: {err}')
                        
                        #sys.exit('The above error is critical, run stopped 2.')
                        continue
                        
        

        img_path = location
        
        return img_path
        
    else:
        
        print(f'Source {source} is not supported for now.')
        
        pass
    
        

   

    
def wait_for_condition(condition_func, timeout=60, check_interval=1):
    start_time = time.time()
    while time.time() - start_time < timeout:
        if condition_func():
            return True
        time.sleep(check_interval)
    return False


def load_bloomberg_addin(excel):
    try:
        
        load_com_addin(excel, "Bloomberg_COM_Addin")        
        #excel.RegisterXLL("C:\\Program Files\\blp\\API\\Office Tools\\BloombergUI.xla")
        excel.Workbooks.Open("C:\\Program Files\\blp\\API\\Office Tools\\BloombergUI.xla")
        excel.RegisterXLL("C:\\Program Files\\blp\\API\\Office Tools\\BloombergUI.xla")  
        
        print("Bloomberg add-in loaded successfully.\n")
    except Exception as e:
        print(f"Failed to load Bloomberg add-in: {e} \n")
        logging.warning(f"Failed to load Bloomberg add-in: {e} \n")
        
        
        
        
def load_com_addin(excel, addin_prog_id):
    try:
        # Access the COM Add-ins
        com_addins = excel.COMAddIns

        # Loop through the add-ins to find the one you want
        for addin in com_addins:
            if addin.ProgID == addin_prog_id:
                if not addin.Connect:  # If the add-in is not connected
                    addin.Connect = True  # Load and connect the add-in
                print(f"{addin.Description} loaded successfully!")
                break
        else:
            print(f"Add-in with ProgID {addin_prog_id} not found.")
    except Exception as e:
        print(f"Failed to load COM add-in '{addin_prog_id}': {e}")
        

def refresh_bloomberg_data(sheet):
    try:
        # Clear any existing errors
        sheet.UsedRange.Errors.Clear()
        print("Cleared existing errors in the sheet")
        
        # Use Bloomberg's RefreshSheet function
        sheet.Parent.Application.Run("RefreshSheet")
        print("Bloomberg RefreshSheet function called")
        
        # Wait for the refresh to complete
        time.sleep(8)  # Adjust this value based on how long the refresh typically takes
        
        # Recalculate the entire sheet
        sheet.Calculate()
        print("Sheet recalculated")
        
        # Additional wait to ensure all data is updated
        time.sleep(5)
        
        print("Bloomberg data refresh completed")
    except Exception as e:
        print(f"Failed to refresh Bloomberg data: {e}")
        logging.warning(f"Failed to refresh Bloomberg data: {e}")
        





def refresh_and_get_screenshot(original_file_path, sheet_name, range_address, output_folder, output_img_name=None, timeout=15, ready_check_range=None, ready_check_sheet=None):
    
    
    opened_files_meta = load_opened_files_meta() # Load already-opened Excel files path data.
    
    # Generate a copy file path for possible use if the file is copied for the 1st time.
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    file_name = os.path.basename(original_file_path)
    base_name, extension = os.path.splitext(file_name)
    copy_file_name = f"{base_name}_copy_{timestamp}{extension}"
    copy_file_path = os.path.join(output_folder, copy_file_name)
    
    
    # Check if the file is already opened
    if original_file_path in opened_files_meta:
        excel = win32com.client.gencache.EnsureDispatch("Excel.Application")  # Ensure Excel is running
        
        excel.Calculation = constants.xlCalculationAutomatic
        
        excel.DisplayAlerts = False # Turn off warnings to force the run move fwd.
        
        excel.Visible = True
        
        
        # Open the copied file instead of the original file
        wb = excel.Workbooks.Open(opened_files_meta[original_file_path]['copied_file_path'])
        

        
    else:

        shutil.copy2(original_file_path, copy_file_path)
        print(f'Created copy of the original file: {copy_file_path} \n')
        logging.info(f'Created copy of the original file: {copy_file_path} \n')
    
        excel = None
        wb = None

        try:
            
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create Excel application without opening a workbook
            excel = win32com.client.gencache.EnsureDispatch("Excel.Application")
            excel.Calculation = constants.xlCalculationAutomatic
            
            excel.DisplayAlerts = False # Turn off warnings to force the run move fwd.
            
            excel.Visible = True
            
            
            
            load_bloomberg_addin(excel)
            
            load_com_addin(excel, 'PowerlinkCOMAddIn.COMAddIn') # load Workspace
            time.sleep(6)
    
    
            def try_open_workbook():
                nonlocal wb
                try:
                    wb = excel.Workbooks.Open(copy_file_path, UpdateLinks=True)
                    
                    wb.Application.AutomationSecurity = 1 # Ensure macros are enabled at security low level.
                    
                    # Check and change to edit mode immediately after opening
                    if wb.ReadOnly:
                        wb.ChangeFileAccess(Mode=3)  # xlReadWrite mode
                        print("Switched workbook to edit mode.\n")
                        
                    _ = wb.Sheets.Count # This will raise an error if the workbook is not fully open.
                    return True
                except Exception as e:
                    print(f"Waiting for file to open... ({e})")
                    return False
    
            if not wait_for_condition(try_open_workbook, timeout):
                
                
                #raise TimeoutError(f"Failed to open the file within {timeout} seconds.\n")
                logging.error(f"Failed to open the file within {timeout} seconds.\n")

            else: # Successful opened.  
            
            
                
                
                opened_files_meta[original_file_path] = {
                'copied_file_path': copy_file_path  # Store the copied file path
                }
                save_opened_files_meta(opened_files_meta) # Update the json file with previously loaded info and currently written info.
            
                
            # Force a recalculation
            excel.CalculateFull()   
                
            time.sleep(10) 
                
           
        
        except Exception as e:
            
            print(f"An error occurred while opening the copied file: {e}")
            
            return
        
    

    if sheet_name not in [sheet.Name for sheet in wb.Sheets]:
        
        #raise ValueError(f"Sheet '{sheet_name}' not found in the workbook.\n")
        logging.error(f"Sheet '{sheet_name}' not found in the workbook.\n")
        
        return

    sheet = wb.Sheets(sheet_name)        
        
    try:    
        # Check if the intended sheet is already selected
        if excel.ActiveSheet.Name != sheet.Name:
            # Ensure sheet is selected and scrolled to top-left
            sheet.Select()
            sheet.Range('A1').Select()
    
        sheet.Range('A1').Select()
                
        # Force another recalculation
        excel.CalculateFull()               
        time.sleep(6)
        
        
        # Refresh Bloomberg data and calculate the sheet
        # refresh_bloomberg_data(sheet) # This function is not working as expected, use pyautogui instead.
        win32gui.SetForegroundWindow(excel.Hwnd)
        win32gui.ShowWindow(excel.Hwnd, win32con.SW_MAXIMIZE)
        
        time.sleep(4)
        pyautogui.hotkey('alt', 'y', '2')
        time.sleep(0.5)
        pyautogui.hotkey('y', '9')
        time.sleep(0.5)
        pyautogui.press('r')
        time.sleep(6)
        
 
        print("Current workbook and sheet recalculation performed.\n")
        logging.info("Current workbook and sheet recalculation performed.\n")
        time.sleep(1)
        
            
        # Use the screenshot range as the default for ready_check_range if not provided
        if ready_check_range is None:
            ready_check_range = range_address
    
        def is_sheet_ready():
            try:
                # Determine the sheet to check
                check_sheet = sheet if ready_check_sheet is None else wb.Sheets(ready_check_sheet)
                check_range = check_sheet.Range(ready_check_range)
                
                # Store the original values of the check range
                original_values = [list(row) for row in check_range.Value]
    
                # Recalculate the sheet
                check_sheet.Calculate()
    
                for i in range(1, check_range.Rows.Count + 1):
                    for j in range(1, check_range.Columns.Count + 1):
                        # Access the cell by row and column
                        cell = check_range.Cells(i, j)
                        
                        # Skip if the cell was originally blank
                        if original_values[i-1][j-1] is None or original_values[i-1][j-1] == "":
                            continue
                        
                        cell_value = cell.Value
                        
                        # Check various conditions that indicate the sheet is not ready
                        if cell_value is None or cell_value == "":
                            return False
                        
                        # Check for NA() formula result
                        if isinstance(cell_value, str) and cell_value.upper().startswith("#N/A"):
                            return False
                        
                        # Check for unrefreshed formulas (starting with '=')
                        if isinstance(cell_value, str) and cell_value.startswith("="):
                            return False
                        
                        # Check for empty strings after stripping whitespace
                        if isinstance(cell_value, str) and not cell_value.strip():
                            return False
                        
                        # Check for #VALUE! error
                        if isinstance(cell_value, str) and cell_value.upper() == "#VALUE!":
                            return False
    
                # If all cells pass, the sheet is ready
                return True
    
            except Exception as e:
                print(f"Error in is_sheet_ready: {e}")
                return False
    
        if not wait_for_condition(is_sheet_ready, timeout):
            # raise TimeoutError(f"Sheet did not become ready within {timeout} seconds.")
            logging.error(f"Sheet did not become ready within {timeout} seconds.")
    
            
            
        # Force a screen update
        excel.ScreenUpdating = True
        time.sleep(2)
        
        
        # Clear clipboard before taking a screenshot    
        ctypes.windll.user32.OpenClipboard(0)
        ctypes.windll.user32.EmptyClipboard()
        ctypes.windll.user32.CloseClipboard()
    
        rng = sheet.Range(range_address)
    
        if output_img_name is None:
            output_img_name = f"{sheet_name}_{range_address}"
    
        screenshot_file = os.path.join(output_folder, f"{output_img_name}.png")
        
        # Try multiple times to capture the screenshot
        max_attempts = 3
        for attempt in range(max_attempts):
            rng.CopyPicture(Appearance=1, Format=2)  # xlScreen, xlBitmap
            
            # Use PIL's ImageGrab to get the image from clipboard
            try:
                img = ImageGrab.grabclipboard()
                if img:
                    img.save(screenshot_file, 'PNG')
                    if os.path.getsize(screenshot_file) > 0:
                        print(f"Screenshot saved: {screenshot_file}")
                        logging.info(f"Screenshot saved: {screenshot_file}")
                        break # Break out of the loop if successful.
                else:
                    print(f"Attempt {attempt + 1} failed. Clipboard does not contain image data. Retrying...")
                    logging.warning(f"Attempt {attempt + 1} failed. Clipboard does not contain image data. Retrying...")
                    
                time.sleep(2)
                    
            except Exception as e:
                print(f"Error saving screenshot on attempt {attempt + 1}: {e}")
                logging.error(f"Error saving screenshot on attempt {attempt + 1}: {e}")
            
            time.sleep(2)
        else:
            print("Failed to capture a non-empty screenshot after multiple attempts.")
            logging.error(f"Failed to capture a non-empty screenshot after multiple attempts.")
            #sys.exit(1)
            
    except Exception as e:
        print(f"An error occurred: {e}")
        #sys.exit(1)
        
    finally:

        if wb is not None:
            del wb
            time.sleep(1)
# =============================================================================
#         if excel:
#             del excel
#             time.sleep(2)
# =============================================================================
            
            
# =============================================================================
#     finally:
#         try:
#             if wb is not None:
# 
#                 wb.Close(SaveChanges=False)  # Close the workbook without saving
#                 time.sleep(6)
# 
#             while True:
#                 try:
#                     # Attempt to access the workbook object
#                     wb_name = wb.Name  # This will raise an error if the workbook is closed
#                     print(f"The workbook '{wb_name}' is still open.")
#                 except Exception as e:
#                     print("The workbook has been closed successfully.")
#                     break  # Exit the loop if the workbook is closed
# 
#                 # Sleep for a specified interval before checking again
#                 time.sleep(1)  # Check every second            
#             
#             
# # =============================================================================
# #             if excel:
# #                 excel.Quit()
# #                 time.sleep(6)
# #                 print(f'File closed: {copy_file_path}')
# # =============================================================================
# 
#         except Exception as e:
#             print(f"Error while closing the workbook: {e}")
#             logging.error(f"Error while closing the workbook: {e}")
#             
#         finally:
#             # Release COM objects
#             if wb is not None:
#                 del wb
#             if excel:
#                 del excel
#                 time.sleep(3)
#             # Uninitialize COM
#             pythoncom.CoUninitialize()
#             # Force garbage collection
#             gc.collect()
#             
#             # Force Excel to quit using os.system (Windows-specific)
#             # os.system("taskkill /f /im excel.exe")
#         
#         
#     
#     #Optionally, remove the copy after processing
#     os.remove(copy_file_path)
#     print(f"Removed temporary copy: {copy_file_path}")
# =============================================================================


def save_opened_files_meta(opened_files_meta):
    with open(OPENED_FILES_META_PATH, "w") as f:
        json.dump(opened_files_meta, f)  # Now it only contains serializable data. Given the previously cumulated info has been loaded into opened_files_meta through load_opened_files_meta(), then this dump in 'w' mode has all the info up to the time.

# Function to load opened files metadata from a JSON file
def load_opened_files_meta():
    # Initialize the opened_files_meta dictionary
    opened_files_meta = {}
    
    # Check if the file exists, if not create it
    if not os.path.exists(OPENED_FILES_META_PATH):
        with open(OPENED_FILES_META_PATH, "w") as f:
            json.dump({}, f)  # Create an empty JSON file
        print(f"Created empty metadata file at: {OPENED_FILES_META_PATH}")
    
    # Load the metadata from the file
    with open(OPENED_FILES_META_PATH, "r") as f:
        opened_files_meta = json.load(f)
    
    return opened_files_meta  # Return the loaded metadata


# Function to close all opened Excel files
def close_all_opened_files(opened_files_meta, remove_copied_files=True):
    # Check if there are any opened files to close
    if not opened_files_meta:
        print("No opened files to close.")
        return  # Exit the function if there are no files

    # Clear clipboard before taking a screenshot    
    ctypes.windll.user32.OpenClipboard(0)
    ctypes.windll.user32.EmptyClipboard()
    ctypes.windll.user32.CloseClipboard()


    # Create a single instance of the Excel application
    excel = win32com.client.gencache.EnsureDispatch("Excel.Application")
    time.sleep(5)
    
    for original_file_path, meta in opened_files_meta.items():
        copied_file_path = meta['copied_file_path']  # Get the copied file path
        
        open_workbook = None
        
        # Check if the copied workbook is already open
        for wb in excel.Workbooks:
            if wb.FullName == copied_file_path:
                open_workbook = wb
                break
        
        if open_workbook is not None:
            # Close the workbook
            open_workbook.Close(SaveChanges=False)  # Change to True if you want to save changes
            time.sleep(3)  # Add a small delay to ensure the file is released
                
            if wb is not None:
                del wb

            while True:
                try:
                    # Attempt to access the workbook object
                    wb_name = wb.Name  # This will raise an error if the workbook is closed
                    print(f"The workbook '{wb_name}' is still open.")
                except Exception as e:
                    print("The workbook has been closed successfully.")
                    break  # Exit the loop if the workbook is closed

                # Sleep for a specified interval before checking again
                time.sleep(1)  # Check every second  
            
        # Check if we need to remove the copied file
        if remove_copied_files:
            # Remove the copied file if it exists
            if os.path.exists(copied_file_path):
                try:
                    os.remove(copied_file_path)
                    print(f"Removed copied file: {copied_file_path}")
                except Exception as e:
                    print(f"Error removing copied file: {e}")
        else:
            print(f"Copied file '{copied_file_path}' is not currently open.")
    
    if excel:
        
        del excel
    
    # Quit the Excel application after processing all files
    #excel.Quit()
    # Uninitialize COM
    pythoncom.CoUninitialize()
    # Force garbage collection
    gc.collect()
    
    opened_files_meta.clear()  # Clear the metadata after closing
    save_opened_files_meta(opened_files_meta)  # Save the updated metadata to the file


    
    
def generate_chart_id(length=10):
    '''
    Generate uppercase alphanumeric string with speicified length.
    '''
    
    x = ''.join(random.choices(string.ascii_uppercase + string.digits, k=length))

    return x   
    
    
    
    

           
    
    